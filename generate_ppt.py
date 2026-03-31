"""
TossPayments 결제경로 제작 PPT 생성 스크립트
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── 색상 상수 ──
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_BG = RGBColor(0x1E, 0x29, 0x3B)
BLUE = RGBColor(0x00, 0x53, 0xD6)
LIGHT_GRAY = RGBColor(0xF1, 0xF5, 0xF9)
GRAY_TEXT = RGBColor(0x64, 0x74, 0x8B)
RED = RGBColor(0xDC, 0x26, 0x26)
TOSS_BLUE = RGBColor(0x00, 0x64, 0xFF)


def add_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_bar(slide, title_text, subtitle_text=""):
    # 상단 타이틀 바
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = TOSS_BLUE
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.6)
    tf.margin_top = Inches(0.15)

    if subtitle_text:
        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.size = Pt(14)
        p2.font.color.rgb = RGBColor(0xCC, 0xDD, 0xFF)
        p2.alignment = PP_ALIGN.LEFT


def add_screenshot_placeholder(slide, left, top, width, height, label):
    """빨간 테두리 스크린샷 영역"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xFF, 0xFA, 0xFA)
    shape.line.color.rgb = RED
    shape.line.width = Pt(2.5)
    shape.line.dash_style = 2  # dash

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    run = p.add_run()
    run.text = "[Screenshot]"
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = RED

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = label
    run2.font.size = Pt(13)
    run2.font.color.rgb = RGBColor(0x99, 0x33, 0x33)


def add_text_box(slide, left, top, width, height, text, font_size=14, bold=False, color=BLACK, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return tf


# ═══════════════════════════════════════════════════════════════════════════════
# 슬라이드 1: 표지
# ═══════════════════════════════════════════════════════════════════════════════
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide1, TOSS_BLUE)

add_text_box(slide1, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
             "PronunFit 결제경로 제작 가이드", font_size=40, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(slide1, Inches(1), Inches(3.0), Inches(11), Inches(0.8),
             "토스페이먼츠 심사용 자료", font_size=24, color=RGBColor(0xCC, 0xDD, 0xFF), alignment=PP_ALIGN.CENTER)

# 하단 정보
info_lines = "상호명: 아리젬스  |  대표: 하승우  |  사업자번호: 746-11-03230\nhttps://multi-translator-seven.vercel.app"
add_text_box(slide1, Inches(1), Inches(5.0), Inches(11), Inches(1),
             info_lines, font_size=16, color=WHITE, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# 슬라이드 2: ① 가맹점 정보 기재
# ═══════════════════════════════════════════════════════════════════════════════
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide2, WHITE)
add_title_bar(slide2, "① 가맹점 정보 기재", "상호명, 사업자등록번호, URL, 테스트 계정 정보")

# 정보 테이블 형식
info_data = [
    ("상호명 (서비스명)", "아리젬스 (PronunFit)"),
    ("대표자명", "하승우"),
    ("사업자등록번호", "746-11-03230"),
    ("서비스 URL", "https://multi-translator-seven.vercel.app"),
    ("이메일", "SystemAdmin@PronunFit.com"),
    ("전화번호", "050-6754-5465"),
    ("사업장 주소", "경기 김포시 걸포2로 83"),
    ("테스트 계정", "Google 소셜 로그인 (별도 계정 불필요)"),
    ("결제 방식", "토스페이먼츠 빌링키 정기결제 (카드)"),
]

y_start = Inches(1.6)
for i, (label, value) in enumerate(info_data):
    y = y_start + Inches(i * 0.55)
    # label bg
    shape = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), y, Inches(3), Inches(0.45))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xEE, 0xF2, 0xFF)
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Inches(0.15)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x33, 0x38, 0xCA)

    # value
    shape2 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), y, Inches(7), Inches(0.45))
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = WHITE
    shape2.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
    shape2.line.width = Pt(1)
    tf2 = shape2.text_frame
    tf2.margin_left = Inches(0.15)
    tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
    p2 = tf2.paragraphs[0]
    p2.text = value
    p2.font.size = Pt(14)
    p2.font.color.rgb = BLACK


# ═══════════════════════════════════════════════════════════════════════════════
# 슬라이드 3: ② 하단정보 캡처 (푸터 사업자 정보)
# ═══════════════════════════════════════════════════════════════════════════════
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide3, WHITE)
add_title_bar(slide3, "② 하단정보 캡처", "랜딩페이지 푸터에 표시되는 사업자 정보")

add_text_box(slide3, Inches(0.8), Inches(1.6), Inches(5), Inches(0.5),
             "랜딩페이지(https://multi-translator-seven.vercel.app) 하단 푸터",
             font_size=13, color=GRAY_TEXT)

add_screenshot_placeholder(slide3, Inches(0.8), Inches(2.2), Inches(5.5), Inches(4.5),
    "랜딩페이지 하단 푸터 캡처\n\n아리젬스 | 대표 하승우\n사업자등록번호: 746-11-03230\n주소, 이메일, 전화번호가\n표시된 영역을 캡처하세요")

# 오른쪽 설명
add_text_box(slide3, Inches(7), Inches(2.2), Inches(5.5), Inches(3),
    "캡처 방법:\n\n"
    "1. https://multi-translator-seven.vercel.app 접속\n\n"
    "2. 페이지 맨 아래로 스크롤\n\n"
    "3. 아래 정보가 보이는 영역 캡처:\n"
    "   • 아리젬스 | 대표 하승우\n"
    "   • 사업자등록번호: 746-11-03230\n"
    "   • 경기 김포시 걸포2로 83\n"
    "   • SystemAdmin@PronunFit.com\n"
    "   • 050-6754-5465\n\n"
    "4. 개인정보처리방침 / 이용약관 / 연락처\n"
    "   링크도 함께 보이도록 캡처",
    font_size=12, color=RGBColor(0x47, 0x55, 0x69))


# ═══════════════════════════════════════════════════════════════════════════════
# 슬라이드 4: ③ 환불규정 캡처 (이용약관 6조)
# ═══════════════════════════════════════════════════════════════════════════════
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide4, WHITE)
add_title_bar(slide4, "③ 환불규정 캡처", "이용약관 내 '제6조 결제 취소 및 환불' 조항")

add_screenshot_placeholder(slide4, Inches(0.8), Inches(2.0), Inches(5.5), Inches(5),
    "이용약관 > 6. 결제 취소 및 환불 캡처\n\n"
    "① 구독 취소\n"
    "② 환불 규정 (1개월/3개월)\n"
    "③ 환불 불가 사유\n"
    "④ 예외적 환불\n\n"
    "전체가 보이도록 캡처하세요")

add_text_box(slide4, Inches(7), Inches(2.0), Inches(5.5), Inches(4),
    "캡처 방법:\n\n"
    "1. 랜딩페이지 하단 '이용약관' 클릭\n\n"
    "2. '6. 결제 취소 및 환불' 항목으로 스크롤\n\n"
    "3. 아래 4개 하위 항목이 모두 보이도록 캡처:\n"
    "   ① 구독 취소 (자동 갱신 중지)\n"
    "   ② 환불 규정 (7일 이내 전액, 일할계산 등)\n"
    "   ③ 환불 불가 사유\n"
    "   ④ 예외적 환불 (결제오류, 장기중단)\n\n"
    "※ 스크롤이 필요하면 여러 장 캡처 가능",
    font_size=12, color=RGBColor(0x47, 0x55, 0x69))


# ═══════════════════════════════════════════════════════════════════════════════
# 슬라이드 5: ④ 로그인/회원가입 캡처
# ═══════════════════════════════════════════════════════════════════════════════
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide5, WHITE)
add_title_bar(slide5, "④ 로그인 / 회원가입 캡처", "Google OAuth 소셜 로그인 과정")

add_screenshot_placeholder(slide5, Inches(0.5), Inches(2.0), Inches(3.8), Inches(4.8),
    "로그인 화면 캡처\n\n"
    "랜딩페이지에서\n'직접 체험해 보세요' 클릭 후\n"
    "Google 로그인 팝업이\n표시된 화면")

add_screenshot_placeholder(slide5, Inches(4.8), Inches(2.0), Inches(3.8), Inches(4.8),
    "Google 계정 선택 캡처\n\n"
    "Google OAuth\n계정 선택 화면")

add_screenshot_placeholder(slide5, Inches(9.1), Inches(2.0), Inches(3.8), Inches(4.8),
    "로그인 완료 캡처\n\n"
    "로그인 성공 후\n메인 번역 화면이\n표시된 상태")


# ═══════════════════════════════════════════════════════════════════════════════
# 슬라이드 6: ⑤ 상품 선택 캡처 (구독 플랜)
# ═══════════════════════════════════════════════════════════════════════════════
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide6, WHITE)
add_title_bar(slide6, "⑤ 상품 선택 / 구매 과정 캡처", "구독 플랜 선택 화면 (UpgradeModal)")

add_screenshot_placeholder(slide6, Inches(0.5), Inches(2.0), Inches(5.8), Inches(5),
    "업그레이드 모달 캡처\n\n"
    "설정 > '구독 업그레이드' 클릭 시\n표시되는 플랜 선택 팝업\n\n"
    "• Pro 1개월 ₩9,900 / 3개월 ₩16,500\n"
    "• Premium 1개월 ₩19,900 / 3개월 ₩55,000")

add_text_box(slide6, Inches(6.8), Inches(2.0), Inches(5.8), Inches(4.5),
    "상품 구성:\n\n"
    "Pro 등급:\n"
    "  • 1개월: ₩9,900 (자동 갱신)\n"
    "  • 3개월: ₩16,500 (일시불, 44% 할인)\n"
    "  • 무제한 발음 평가\n"
    "  • AI 번역 팁 무제한\n\n"
    "Premium 등급:\n"
    "  • 1개월: ₩19,900 (자동 갱신)\n"
    "  • 3개월: ₩55,000 (일시불, 8% 할인)\n"
    "  • Pro 기능 전부 포함\n"
    "  • AI 문법/뉘앙스 분석\n"
    "  • 고급 발음 코칭\n\n"
    "결제 방식: 토스페이먼츠 빌링키 (카드)",
    font_size=12, color=RGBColor(0x47, 0x55, 0x69))


# ═══════════════════════════════════════════════════════════════════════════════
# 슬라이드 7: ⑤ 상품 상세 정보
# ═══════════════════════════════════════════════════════════════════════════════
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide7, WHITE)
add_title_bar(slide7, "⑤ 상품 상세 정보", "각 구독 플랜 기능 비교")

# 테이블 형식 - 헤더
headers = ["기능", "Free Trial", "Pro", "Premium"]
col_widths = [Inches(4), Inches(2.5), Inches(2.5), Inches(2.5)]
x_start = Inches(0.8)
y_header = Inches(1.8)

x = x_start
for i, (header, w) in enumerate(zip(headers, col_widths)):
    shape = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y_header, w, Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = TOSS_BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = header
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    x += w

# 테이블 행
rows = [
    ("번역 (텍스트/음성/카메라)", "하루 5회", "무제한", "무제한"),
    ("발음 평가", "하루 3회", "무제한", "무제한"),
    ("AI 번역 팁", "하루 3회", "무제한", "무제한"),
    ("보관함 저장", "최대 20개", "무제한", "무제한"),
    ("AI 문법/뉘앙스 분석", "✕", "✕", "무제한"),
    ("고급 발음 코칭", "✕", "✕", "무제한"),
    ("광고 제거", "✕", "✕", "✓"),
]

for ri, (feat, trial, pro, premium) in enumerate(rows):
    y = y_header + Inches(0.5) + Inches(ri * 0.5)
    bg = WHITE if ri % 2 == 0 else LIGHT_GRAY
    vals = [feat, trial, pro, premium]
    x = x_start
    for ci, (val, w) in enumerate(zip(vals, col_widths)):
        shape = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Inches(0.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg
        shape.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
        shape.line.width = Pt(0.5)
        tf = shape.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.1)
        p = tf.paragraphs[0]
        p.text = val
        p.font.size = Pt(12)
        p.font.color.rgb = BLACK if ci == 0 else (RGBColor(0x47, 0x55, 0x69))
        p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
        x += w


# ═══════════════════════════════════════════════════════════════════════════════
# 슬라이드 8: ⑥ 카드 결제 경로 캡처
# ═══════════════════════════════════════════════════════════════════════════════
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide8, WHITE)
add_title_bar(slide8, "⑥ 카드 결제 경로 캡처", "토스페이먼츠 카드 등록 → 결제 완료 과정")

add_screenshot_placeholder(slide8, Inches(0.3), Inches(2.0), Inches(3.1), Inches(5),
    "플랜 선택 후\n'시작하기' 버튼 클릭\n\n"
    "UpgradeModal에서\n원하는 플랜의\n'시작하기' 클릭")

add_screenshot_placeholder(slide8, Inches(3.7), Inches(2.0), Inches(3.1), Inches(5),
    "토스 카드 등록 화면\n\n"
    "토스페이먼츠\n빌링키 인증 화면\n(카드번호 입력)")

add_screenshot_placeholder(slide8, Inches(7.1), Inches(2.0), Inches(3.1), Inches(5),
    "결제 완료 화면\n\n"
    "카드 등록 성공 후\n결제 완료 및\n구독 활성화 화면")

add_text_box(slide8, Inches(10.5), Inches(2.0), Inches(2.5), Inches(5),
    "결제 흐름:\n\n"
    "1. 플랜 선택\n   (시작하기 클릭)\n\n"
    "2. 토스 카드등록\n   (빌링키 인증)\n\n"
    "3. 자동 결제 실행\n   (서버에서 처리)\n\n"
    "4. 구독 활성화\n   (Firestore 업데이트)\n\n"
    "5. 완료 페이지\n   (성공 메시지)",
    font_size=11, color=RGBColor(0x47, 0x55, 0x69))


# ═══════════════════════════════════════════════════════════════════════════════
# 슬라이드 9: 결제 흐름도 (텍스트 기반)
# ═══════════════════════════════════════════════════════════════════════════════
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide9, WHITE)
add_title_bar(slide9, "결제 시스템 흐름도", "빌링키 기반 정기결제 아키텍처")

flow_steps = [
    ("1. 플랜 선택", "사용자가 Pro/Premium\n1개월 또는 3개월 선택", TOSS_BLUE),
    ("2. 빌링키 인증", "토스페이먼츠 SDK\nrequestBillingAuth()\n카드 정보 입력", RGBColor(0x43, 0x38, 0xCA)),
    ("3. 콜백 처리", "successUrl로 리다이렉트\nauthKey 수신", RGBColor(0x05, 0x96, 0x69)),
    ("4. 빌링키 발급", "서버에서 authKey →\nbillingKey 교환\nFirestore 저장", RGBColor(0xB4, 0x53, 0x09)),
    ("5. 결제 실행", "billingKey로 즉시 결제\n토스 confirm-billing API", RGBColor(0xDC, 0x26, 0x26)),
    ("6. 구독 활성화", "Firestore tier 업데이트\n만료일 설정\nautoRenew: true", RGBColor(0x05, 0x96, 0x69)),
]

for i, (title, desc, color) in enumerate(flow_steps):
    x = Inches(0.4) + Inches(i * 2.15)
    y = Inches(2.5)

    # 번호 원
    shape = slide9.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.7), Inches(1.8), Inches(0.5), Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = str(i + 1)
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # 박스
    shape2 = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2), Inches(2.5))
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = WHITE
    shape2.line.color.rgb = color
    shape2.line.width = Pt(2)

    tf2 = shape2.text_frame
    tf2.word_wrap = True
    tf2.margin_left = Inches(0.1)
    tf2.margin_right = Inches(0.1)
    tf2.margin_top = Inches(0.15)

    p_title = tf2.paragraphs[0]
    p_title.text = title
    p_title.font.size = Pt(13)
    p_title.font.bold = True
    p_title.font.color.rgb = color
    p_title.alignment = PP_ALIGN.CENTER

    p_desc = tf2.add_paragraph()
    p_desc.text = "\n" + desc
    p_desc.font.size = Pt(11)
    p_desc.font.color.rgb = GRAY_TEXT
    p_desc.alignment = PP_ALIGN.CENTER

    # 화살표 (마지막 제외)
    if i < len(flow_steps) - 1:
        add_text_box(slide9, x + Inches(2), Inches(3.4), Inches(0.2), Inches(0.4),
                     "→", font_size=20, bold=True, color=RGBColor(0x94, 0xA3, 0xB8), alignment=PP_ALIGN.CENTER)

# 자동갱신 설명
add_text_box(slide9, Inches(0.5), Inches(5.5), Inches(12), Inches(1.5),
    "※ 1개월 플랜: 매월 자동 갱신 (서버 cron job이 만료 3일 전 자동 결제)\n"
    "※ 3개월 플랜: 일시불 결제 (자동 갱신 없음, 만료 시 Free Trial로 전환)\n"
    "※ 자동 갱신 중지: 설정 > 업그레이드 모달 > '자동 갱신 중지' 버튼",
    font_size=12, color=GRAY_TEXT)


# ═══════════════════════════════════════════════════════════════════════════════
# 슬라이드 10: 마무리
# ═══════════════════════════════════════════════════════════════════════════════
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide10, TOSS_BLUE)

add_text_box(slide10, Inches(1), Inches(2.0), Inches(11), Inches(1),
             "PronunFit 결제경로 제작 가이드", font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_text_box(slide10, Inches(1), Inches(3.5), Inches(11), Inches(2),
    "상호명: 아리젬스 (PronunFit)\n"
    "대표: 하승우 | 사업자번호: 746-11-03230\n"
    "URL: https://multi-translator-seven.vercel.app\n"
    "이메일: SystemAdmin@PronunFit.com | 전화: 050-6754-5465",
    font_size=16, color=WHITE, alignment=PP_ALIGN.CENTER)

add_text_box(slide10, Inches(1), Inches(5.8), Inches(11), Inches(0.5),
    "감사합니다", font_size=24, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# 저장
# ═══════════════════════════════════════════════════════════════════════════════
output_path = r"C:\Users\User\Desktop\PronunFit_결제경로_제작파일.pptx"
prs.save(output_path)
print(f"PPT saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
